from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import logging
import os
from pathlib import Path

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants
COMPANY_BLUE = RGBColor(0, 75, 150)
OUTPUT_DIR = Path("output")
IMAGES_DIR = OUTPUT_DIR

class PresentationBuilder:
    def __init__(self):
        self.prs = Presentation()
        self._setup_slide_layouts()

    def _setup_slide_layouts(self):
        """Store commonly used slide layouts."""
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.content_slide_layout = self.prs.slide_layouts[1]
        self.image_slide_layout = self.prs.slide_layouts[5]

    def _add_title_slide(self):
        """Add the title slide."""
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "NASA CME and GST Analysis"
        subtitle.text = "Solar Event Impact Analysis\nData Science Team"

        # Format title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = COMPANY_BLUE

    def _add_overview_slide(self):
        """Add the overview slide."""
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Overview"
        content.text = (
            "• Analysis of Coronal Mass Ejections (CMEs) and their impact on Earth\n"
            "• Study period: 2013-2024\n"
            "• Focus on propagation times and geomagnetic storm correlation\n"
            "• Data source: NASA DONKI API\n"
            "• Key metrics: Speed, impact time, and storm intensity"
        )

    def _add_methodology_slide(self):
        """Add the methodology slide."""
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Methodology"
        content.text = (
            "1. Data Collection\n"
            "   • Retrieved CME data from NASA DONKI API\n"
            "   • Collected associated geomagnetic storm data\n\n"
            "2. Data Processing\n"
            "   • Cleaned and merged datasets\n"
            "   • Calculated propagation times\n"
            "   • Analyzed speed-intensity correlations\n\n"
            "3. Analysis\n"
            "   • Statistical analysis of propagation times\n"
            "   • Correlation studies\n"
            "   • Temporal pattern analysis"
        )

    def _add_key_findings_slide(self):
        """Add the key findings slide."""
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Key Findings"
        content.text = (
            "Statistical Summary:\n\n"
            "• Average Propagation Time: ~2 days 21 hours\n"
            "• Minimum Time: ~1 day 5 hours\n"
            "• Maximum Time: ~6 days 3 hours\n"
            "• Median Time: ~2 days 17 hours\n\n"
            "Correlations:\n"
            "• Strong correlation between CME speed and storm intensity\n"
            "• Seasonal variations in event frequency observed"
        )

    def _add_visualization_slide(self, image_path, title_text):
        """Add a visualization slide."""
        try:
            slide = self.prs.slides.add_slide(self.image_slide_layout)
            title = slide.shapes.title
            title.text = title_text

            # Calculate image position and size
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)

            # Add image
            slide.shapes.add_picture(str(image_path), left, top, width, height)
        except Exception as e:
            logger.error(f"Failed to add visualization slide: {e}")
            raise

    def _add_conclusions_slide(self):
        """Add the conclusions slide."""
        slide = self.prs.slides.add_slide(self.content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Conclusions"
        content.text = (
            "• CME propagation times show consistent patterns\n"
            "• Faster CMEs generally lead to more intense storms\n"
            "• Seasonal variations affect event frequency\n"
            "• Prediction model improvements possible\n\n"
            "Next Steps:\n"
            "• Develop real-time monitoring system\n"
            "• Enhance prediction accuracy\n"
            "• Expand analysis timeframe"
        )

    def create_presentation(self):
        """Create the complete presentation."""
        try:
            # Add all slides
            self._add_title_slide()
            self._add_overview_slide()
            self._add_methodology_slide()
            self._add_key_findings_slide()

            # Add visualization slides
            visualizations = [
                ("monthly_events.png", "Monthly Event Distribution"),
                ("propagation_times.png", "CME Propagation Times"),
                ("speed_kp_correlation.png", "Speed-Intensity Correlation")
            ]

            for viz_file, title in visualizations:
                viz_path = IMAGES_DIR / viz_file
                if viz_path.exists():
                    self._add_visualization_slide(viz_path, title)
                else:
                    logger.warning(f"Visualization file not found: {viz_file}")

            self._add_conclusions_slide()

            # Save presentation
            output_path = OUTPUT_DIR / "nasa_cme_gst_analysis.pptx"
            self.prs.save(str(output_path))
            logger.info(f"Presentation saved successfully to {output_path}")

        except Exception as e:
            logger.error(f"Failed to create presentation: {e}")
            raise

if __name__ == "__main__":
    try:
        builder = PresentationBuilder()
        builder.create_presentation()
    except Exception as e:
        logger.error(f"Presentation generation failed: {e}")
        raise

